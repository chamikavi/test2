from fastapi import FastAPI, Depends, HTTPException
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from .database import SessionLocal, engine, Base
from . import models, schemas
from passlib.hash import bcrypt
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import tempfile

security = HTTPBasic()
app = FastAPI(title="360 Performance Hub API")

Base.metadata.create_all(bind=engine)

# Dependency
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def get_current_user(credentials: HTTPBasicCredentials = Depends(security), db: Session = Depends(get_db)):
    user = db.query(models.User).filter(models.User.username == credentials.username).first()
    if not user or not bcrypt.verify(credentials.password, user.password):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return user

@app.post("/users/", response_model=schemas.User)
def create_user(
    user: schemas.UserCreate,
    db: Session = Depends(get_db),
    current: models.User = Depends(get_current_user),
):
    if current.role != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    hashed_pw = bcrypt.hash(user.password)
    db_user = models.User(username=user.username, password=hashed_pw, role=user.role)
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    return db_user

@app.get("/outlets/", response_model=list[schemas.Outlet])
def list_outlets(db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.Outlet).all()

@app.post("/outlets/", response_model=schemas.Outlet)
def create_outlet(outlet: schemas.OutletCreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_outlet = models.Outlet(name=outlet.name, manager_id=current.id)
    db.add(db_outlet)
    db.commit()
    db.refresh(db_outlet)
    return db_outlet


@app.post("/periods/", response_model=schemas.Period)
def create_period(period: schemas.PeriodCreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_period = models.Period(month=period.month, year=period.year)
    db.add(db_period)
    db.commit()
    db.refresh(db_period)
    return db_period


@app.get("/periods/", response_model=list[schemas.Period])
def list_periods(db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.Period).order_by(models.Period.year, models.Period.month).all()


@app.post("/kpis/", response_model=schemas.KPI)
def create_kpi(kpi: schemas.KPICreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_kpi = models.KPI(name=kpi.name)
    db.add(db_kpi)
    db.commit()
    db.refresh(db_kpi)
    return db_kpi


@app.get("/kpis/", response_model=list[schemas.KPI])
def list_kpis(db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.KPI).all()


@app.post("/updates/", response_model=schemas.Update)
def create_update(update: schemas.UpdateCreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_update = models.Update(**update.dict())
    db.add(db_update)
    db.commit()
    db.refresh(db_update)
    return db_update


@app.get("/updates/{outlet_id}/{kpi_id}", response_model=list[schemas.Update])
def list_updates(outlet_id: int, kpi_id: int, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.Update).filter(models.Update.outlet_id == outlet_id, models.Update.kpi_id == kpi_id).order_by(models.Update.period_id).all()


@app.get("/metrics/{outlet_id}/{kpi_id}")
def metrics(outlet_id: int, kpi_id: int, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    updates = (
        db.query(models.Update, models.Period)
        .join(models.Period, models.Update.period_id == models.Period.id)
        .filter(models.Update.outlet_id == outlet_id, models.Update.kpi_id == kpi_id)
        .order_by(models.Period.year, models.Period.month)
        .all()
    )
    return [
        {
            "period": f"{period.year}-{period.month:02d}",
            "value": update.value,
            "note": update.note,
        }
        for update, period in updates
    ]


@app.post("/feedback/", response_model=schemas.Feedback)
def create_feedback(feedback: schemas.FeedbackCreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_fb = models.Feedback(**feedback.dict())
    db.add(db_fb)
    db.commit()
    db.refresh(db_fb)
    return db_fb


@app.get("/feedback/{outlet_id}", response_model=list[schemas.Feedback])
def list_feedback(outlet_id: int, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.Feedback).filter(models.Feedback.outlet_id == outlet_id).order_by(models.Feedback.period_id).all()


@app.post("/files/", response_model=schemas.File)
def create_file(file: schemas.FileCreate, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    db_file = models.File(**file.dict())
    db.add(db_file)
    db.commit()
    db.refresh(db_file)
    return db_file


@app.get("/files/{outlet_id}/{period_id}", response_model=list[schemas.File])
def list_files(outlet_id: int, period_id: int, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    return db.query(models.File).filter(models.File.outlet_id == outlet_id, models.File.period_id == period_id).all()


@app.get("/deck/{outlet_id}/{period_id}")
def generate_deck(outlet_id: int, period_id: int, db: Session = Depends(get_db), current: models.User = Depends(get_current_user)):
    period = db.query(models.Period).filter(models.Period.id == period_id).first()
    if not period:
        raise HTTPException(status_code=404, detail="Period not found")
    updates = (
        db.query(models.Update, models.KPI)
        .join(models.KPI, models.Update.kpi_id == models.KPI.id)
        .filter(models.Update.outlet_id == outlet_id, models.Update.period_id == period_id)
        .all()
    )
    if not updates:
        raise HTTPException(status_code=404, detail="No data for deck")

    kpi_names = [kpi.name for _, kpi in updates]
    values = [upd.value for upd, _ in updates]

    fig, ax = plt.subplots()
    ax.bar(kpi_names, values)
    ax.set_ylabel("Value")
    plt.xticks(rotation=45, ha="right")
    tmp_chart = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    plt.tight_layout()
    plt.savefig(tmp_chart.name)
    plt.close(fig)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Outlet {outlet_id} {period.year}-{period.month:02d}"
    slide.shapes.add_picture(tmp_chart.name, Inches(1), Inches(2), height=Inches(4))

    tmp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_pptx.name)

    return FileResponse(tmp_pptx.name, filename=f"deck_{outlet_id}_{period.year}-{period.month:02d}.pptx")
